from __future__ import annotations

import json
import os
import shutil
from decimal import Decimal
from html import escape as html_escape
from pathlib import Path
from typing import Any
from urllib.error import HTTPError, URLError
from urllib.parse import urljoin, urlparse
from urllib.request import Request, urlopen
from uuid import uuid4

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from packages.shared_contracts.python_models import ArtifactDTO, DeckDTO, SlideDraftDTO


ARTIFACT_ROOT = Path("runtime_data/artifacts")
FONT_NAME = "Microsoft YaHei"
COLORS = {
    "bg": RGBColor(10, 27, 44),
    "panel": RGBColor(14, 40, 63),
    "panel_2": RGBColor(18, 52, 79),
    "accent": RGBColor(35, 129, 239),
    "cyan": RGBColor(45, 212, 191),
    "text": RGBColor(235, 244, 255),
    "muted": RGBColor(161, 180, 204),
    "line": RGBColor(60, 96, 132),
    "white": RGBColor(255, 255, 255),
}


def _inches(value: float):
    return Inches(value)


def _safe_text(value: object, fallback: str = "-") -> str:
    text = str(value if value is not None else fallback).strip()
    return text or fallback


def _to_float(value: object) -> float | None:
    if value is None or isinstance(value, bool):
        return None
    if isinstance(value, (int, float, Decimal)):
        return float(value)
    try:
        return float(str(value).replace(",", "").strip())
    except (TypeError, ValueError):
        return None


def _format_cell(value: object) -> str:
    number = _to_float(value)
    if number is None:
        return _safe_text(value, "")
    if abs(number) >= 1000:
        return f"{number:,.0f}"
    if number.is_integer():
        return f"{number:.0f}"
    return f"{number:.2f}"


def _set_run_font(run, *, size: float = 12, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _set_cell_text(cell, text: str, *, size: float = 8, bold: bool = False, color: RGBColor | None = None) -> None:
    cell.text = ""
    frame = cell.text_frame
    frame.margin_left = _inches(0.04)
    frame.margin_right = _inches(0.04)
    frame.margin_top = _inches(0.03)
    frame.margin_bottom = _inches(0.03)
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    _set_run_font(run, size=size, bold=bold, color=color or COLORS["text"])


def _fill_shape(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def _line_shape(shape, color: RGBColor = COLORS["line"], width: float = 1) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def _add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 12,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN | None = None,
) -> Any:
    shape = slide.shapes.add_textbox(_inches(x), _inches(y), _inches(w), _inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = _inches(0.03)
    frame.margin_right = _inches(0.03)
    frame.margin_top = _inches(0.02)
    frame.margin_bottom = _inches(0.02)
    frame.word_wrap = True
    for index, line in enumerate(text.splitlines() or [""]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        if align is not None:
            paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        _set_run_font(run, size=size, bold=bold, color=color or COLORS["text"])
    return shape


def _add_panel(slide, x: float, y: float, w: float, h: float, *, color: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        _inches(x),
        _inches(y),
        _inches(w),
        _inches(h),
    )
    _fill_shape(shape, color or COLORS["panel"], transparency=2)
    _line_shape(shape, COLORS["line"], 0.8)
    return shape


def _add_section_title(slide, x: float, y: float, text: str) -> None:
    _add_textbox(slide, x, y, 4.8, 0.25, text, size=10.5, bold=True, color=COLORS["cyan"])


def _metric_items(slide: SlideDraftDTO) -> list[dict[str, str]]:
    items = slide.lineage_summary.get("key_metrics", [])
    if not isinstance(items, list):
        return []
    result: list[dict[str, str]] = []
    for item in items[:5]:
        if not isinstance(item, dict):
            continue
        result.append({"label": _safe_text(item.get("label")), "value": _safe_text(item.get("value"))})
    return result


def _rows(slide: SlideDraftDTO) -> list[dict[str, Any]]:
    rows = slide.chart_spec.get("rows", [])
    return rows if isinstance(rows, list) else []


def _columns(rows: list[dict[str, Any]], max_columns: int = 5) -> list[str]:
    if not rows or not isinstance(rows[0], dict):
        return []
    return [str(column) for column in rows[0].keys()][:max_columns]


def _chart_fields(rows: list[dict[str, Any]]) -> tuple[list[str], str | None, list[float]]:
    columns = _columns(rows, max_columns=6)
    if not rows or not columns:
        return [], None, []
    numeric_columns = [
        column
        for column in columns
        if any(_to_float(row.get(column)) is not None for row in rows[:8] if isinstance(row, dict))
    ]
    value_column = numeric_columns[0] if numeric_columns else None
    if value_column is None:
        return [], None, []
    category_columns = [column for column in columns if column != value_column][:2]
    labels: list[str] = []
    values: list[float] = []
    for index, row in enumerate(rows[:8]):
        if not isinstance(row, dict):
            continue
        value = _to_float(row.get(value_column))
        if value is None:
            continue
        if category_columns:
            label = " / ".join(_safe_text(row.get(column)) for column in category_columns)
        else:
            label = f"Row {index + 1}"
        labels.append(label[:32])
        values.append(value)
    return labels, value_column, values


def _compact_text(value: object, *, max_chars: int, fallback: str = "-") -> str:
    text = _safe_text(value, fallback).replace("\n", " ").strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


def _schema_text(value: object, *, min_chars: int, max_chars: int, fallback: str) -> str:
    text = _compact_text(value, max_chars=max_chars, fallback=fallback)
    while len(text) < min_chars:
        text = f"{text}。"
    return _compact_text(text, max_chars=max_chars, fallback=fallback)


def _numeric_columns(rows: list[dict[str, Any]], max_columns: int = 6) -> list[str]:
    result: list[str] = []
    for column in _columns(rows, max_columns=max_columns):
        if any(_to_float(row.get(column)) is not None for row in rows if isinstance(row, dict)):
            result.append(column)
    return result


def _text_columns(rows: list[dict[str, Any]], max_columns: int = 6) -> list[str]:
    numeric = set(_numeric_columns(rows, max_columns=max_columns))
    return [column for column in _columns(rows, max_columns=max_columns) if column not in numeric]


def _chart_points(
    rows: list[dict[str, Any]],
    category_column: str,
    value_column: str,
    *,
    max_points: int = 5,
) -> list[dict[str, Any]]:
    totals: dict[str, float] = {}
    counts: dict[str, int] = {}
    for row in rows:
        if not isinstance(row, dict):
            continue
        label = _safe_text(row.get(category_column), "")
        value = _to_float(row.get(value_column))
        if not label or value is None:
            continue
        totals[label] = totals.get(label, 0.0) + value
        counts[label] = counts.get(label, 0) + 1

    points = [
        {
            "name": _compact_text(label, max_chars=18),
            "value": round(totals[label] / max(1, counts[label]), 2),
        }
        for label in totals
    ]
    points.sort(key=lambda item: item["value"], reverse=True)
    return points[:max_points]


def _presenton_icon(query: str) -> dict[str, str]:
    return {"__icon_query__": _compact_text(query, max_chars=20, fallback="chart analysis")}


def _chart_bullets(
    *,
    category_column: str,
    value_column: str,
    points: list[dict[str, Any]],
    slide: SlideDraftDTO,
) -> list[dict[str, Any]]:
    top = points[0] if points else {"name": "-", "value": 0}
    bottom = points[-1] if points else {"name": "-", "value": 0}
    return [
        {
            "title": _schema_text(f"{top['name']}最高", min_chars=2, max_chars=80, fallback="最高项"),
            "description": _schema_text(
                f"{category_column}维度中，{top['name']}的{value_column}最高，数值为{_format_cell(top['value'])}。",
                min_chars=10,
                max_chars=150,
                fallback="最高项表现突出。",
            ),
            "icon": _presenton_icon("trending up chart"),
        },
        {
            "title": _schema_text(f"{bottom['name']}最低", min_chars=2, max_chars=80, fallback="最低项"),
            "description": _schema_text(
                f"{bottom['name']}处于当前图表低位，可作为后续下钻比较对象。",
                min_chars=10,
                max_chars=150,
                fallback="低位项需要继续复核。",
            ),
            "icon": _presenton_icon("compare data"),
        },
        {
            "title": "来源可追溯",
            "description": _schema_text(
                f"本页数据来自 query={slide.query_id or '-'} 的 SQL 结果聚合，保留原始字段用于复核。",
                min_chars=10,
                max_chars=150,
                fallback="本页数据来自 SQL 结果聚合。",
            ),
            "icon": _presenton_icon("database source"),
        },
    ]


def _presenton_direct_markdown(layout_id: str, title: str, summary: str, content: dict[str, Any]) -> str:
    return "\n\n".join(
        [
            f"PRESENTON_LAYOUT_ID: {layout_id}",
            f"# {title}",
            summary,
            "PRESENTON_DIRECT_CONTENT_JSON_START",
            json.dumps(content, ensure_ascii=False),
            "PRESENTON_DIRECT_CONTENT_JSON_END",
        ]
    )


def _presenton_chart_slide(
    slide: SlideDraftDTO,
    *,
    category_column: str,
    value_column: str,
    chart_type: str,
    points: list[dict[str, Any]],
) -> str:
    title = _schema_text(f"{category_column} {value_column}对比", min_chars=3, max_chars=40, fallback=slide.title)
    description = _schema_text(
        f"基于 SQL 结果按{category_column}聚合{value_column}，展示 Top {len(points)} 的差异。",
        min_chars=10,
        max_chars=150,
        fallback="基于 SQL 结果聚合生成图表。",
    )
    content = {
        "title": title,
        "description": description,
        "chartData": {
            "type": chart_type if chart_type in {"bar", "pie", "line", "area", "scatter"} else "bar",
            "data": points,
        },
        "showLegend": False,
        "showTooltip": True,
        "bulletPoints": _chart_bullets(
            category_column=category_column,
            value_column=value_column,
            points=points,
            slide=slide,
        ),
        "__speaker_note__": _schema_text(
            f"本页展示{category_column}维度下{value_column}的聚合对比。"
            f"最高项为{points[0]['name']}，数值为{_format_cell(points[0]['value'])}。"
            "汇报时应先说明口径来自 SQL 查询结果，再解释高低差异及后续复核方向。",
            min_chars=100,
            max_chars=500,
            fallback="本页展示 SQL 数据聚合后的图表结论。",
        ),
    }
    return _presenton_direct_markdown(
        "general:chart-with-bullets-slide",
        title,
        f"图表类型：{content['chartData']['type']}；数据点：{json.dumps(points, ensure_ascii=False)}",
        content,
    )


def _presenton_table_slide(slide: SlideDraftDTO, rows: list[dict[str, Any]]) -> str:
    columns = _columns(rows, max_columns=5)
    numeric_columns = _numeric_columns(rows, max_columns=5)
    sort_column = numeric_columns[0] if numeric_columns else None
    visible_rows = [row for row in rows if isinstance(row, dict)]
    if sort_column:
        visible_rows.sort(key=lambda row: _to_float(row.get(sort_column)) or float("-inf"), reverse=True)
    table_rows = [
        [_compact_text(_format_cell(row.get(column)), max_chars=50, fallback="-") for column in columns]
        for row in visible_rows[:6]
    ]
    while len(table_rows) < 2:
        table_rows.append(["-" for _ in columns])
    title = _schema_text("SQL 结果 Top 明细", min_chars=3, max_chars=40, fallback="SQL 结果明细")
    content = {
        "title": title,
        "tableData": {
            "headers": [_compact_text(column, max_chars=30) for column in columns],
            "rows": table_rows,
        },
        "description": _schema_text(
            f"展示 SQL 结果中最关键的 {len(table_rows)} 行数据，字段包括 {'、'.join(columns)}。",
            min_chars=10,
            max_chars=200,
            fallback="展示 SQL 结果中的关键明细数据。",
        ),
        "__speaker_note__": _schema_text(
            "本页用于承接图表结论，展示支撑结论的明细行。汇报时先说明排序口径，"
            "再指出这些明细与前序图表中最高或最低项目之间的关系，便于业务侧复核数据来源。",
            min_chars=100,
            max_chars=500,
            fallback="本页展示 SQL 结果明细。",
        ),
    }
    return _presenton_direct_markdown(
        "general:table-info-slide",
        title,
        "结构化表格数据来自 SQL 查询结果，交由 Presenton 表格模板渲染。",
        content,
    )


def _short_metric_value(value: object) -> str:
    number = _to_float(value)
    if number is None:
        return _compact_text(value, max_chars=10, fallback="-")
    if abs(number) >= 10000:
        return _compact_text(f"{number / 10000:.1f}万", max_chars=10)
    if abs(number) >= 1000:
        return _compact_text(f"{number:,.0f}", max_chars=10)
    if number.is_integer():
        return str(int(number))
    return _compact_text(f"{number:.2f}", max_chars=10)


def _presenton_metric_items(slide: SlideDraftDTO, rows: list[dict[str, Any]]) -> list[dict[str, str]]:
    metrics: list[dict[str, str]] = []
    raw_metrics = slide.lineage_summary.get("key_metrics", [])
    if isinstance(raw_metrics, list):
        for item in raw_metrics:
            if not isinstance(item, dict):
                continue
            metrics.append(
                {
                    "label": _schema_text(item.get("label"), min_chars=2, max_chars=50, fallback="指标"),
                    "value": _short_metric_value(item.get("value")),
                    "description": _schema_text(
                        f"{item.get('label', '指标')} 当前值为 {item.get('value', '-')}，用于概览本次 SQL 分析范围。",
                        min_chars=10,
                        max_chars=150,
                        fallback="该指标用于概览本次 SQL 分析范围。",
                    ),
                }
            )
            if len(metrics) >= 3:
                return metrics

    columns = _columns(rows, max_columns=6)
    numeric_columns = _numeric_columns(rows, max_columns=6)
    metrics.append(
        {
            "label": "返回行数",
            "value": _short_metric_value(len(rows)),
            "description": _schema_text("本次 SQL 查询返回的原始结果行数，决定图表聚合和明细展示范围。", min_chars=10, max_chars=150, fallback="SQL 返回行数。"),
        }
    )
    metrics.append(
        {
            "label": "字段数",
            "value": _short_metric_value(len(columns)),
            "description": _schema_text("参与本次汇报生成的数据字段数量，用于判断维度和指标覆盖情况。", min_chars=10, max_chars=150, fallback="SQL 字段数量。"),
        }
    )
    if numeric_columns:
        value_column = numeric_columns[0]
        values = [_to_float(row.get(value_column)) for row in rows if isinstance(row, dict)]
        values = [value for value in values if value is not None]
        if values:
            metrics.append(
                {
                    "label": f"{_compact_text(value_column, max_chars=44)}最大",
                    "value": _short_metric_value(max(values)),
                    "description": _schema_text(f"{value_column} 的最大值用于定位价格或指标高点。", min_chars=10, max_chars=150, fallback="最大值用于定位高点。"),
                }
            )
    return metrics[:3]


def _presenton_metrics_slide(slide: SlideDraftDTO, rows: list[dict[str, Any]]) -> str:
    title = _schema_text("关键指标概览", min_chars=3, max_chars=100, fallback="指标概览")
    content = {
        "title": title,
        "metrics": _presenton_metric_items(slide, rows),
        "__speaker_note__": _schema_text(
            "本页作为数据汇报的指标总览，用于快速说明查询规模、字段覆盖和关键数值。"
            "讲解时先给出指标口径，再连接前面图表页中的品牌或场景差异。",
            min_chars=100,
            max_chars=500,
            fallback="本页展示关键指标。",
        ),
    }
    return _presenton_direct_markdown(
        "general:metrics-slide",
        title,
        "关键指标以 Presenton metrics 模板渲染。",
        content,
    )


def _build_presenton_direct_slides(deck: DeckDTO, slides: list[SlideDraftDTO]) -> list[str]:
    max_slides = int(os.getenv("PRESENTON_MAX_DIRECT_SLIDES", "6"))
    result: list[str] = []
    effective_slides = slides or [_placeholder_slide(deck)]
    for slide in effective_slides:
        rows = _rows(slide)
        numeric_columns = _numeric_columns(rows)
        text_columns = _text_columns(rows)
        chart_type = str(slide.chart_spec.get("chart_type") or "bar").lower()
        if chart_type not in {"bar", "pie", "line", "area", "scatter"}:
            chart_type = "bar"

        if rows and numeric_columns and text_columns:
            value_column = numeric_columns[0]
            for category_column in text_columns[:2]:
                points = _chart_points(rows, category_column, value_column)
                if len(points) >= 2:
                    result.append(
                        _presenton_chart_slide(
                            slide,
                            category_column=category_column,
                            value_column=value_column,
                            chart_type=chart_type,
                            points=points,
                        )
                    )
                if len(result) >= max_slides:
                    return result

        if rows and _columns(rows):
            result.append(_presenton_table_slide(slide, rows))
            if len(result) >= max_slides:
                return result

        result.append(_presenton_metrics_slide(slide, rows))
        if len(result) >= max_slides:
            return result

    return result or [
        _presenton_direct_markdown(
            "general:metrics-slide",
            _schema_text(deck.title, min_chars=3, max_chars=100, fallback="数据分析汇报"),
            "暂无 SQL 明细，生成占位指标页。",
            {
                "title": _schema_text(deck.title, min_chars=3, max_chars=100, fallback="数据分析汇报"),
                "metrics": [
                    {"label": "Slide 数", "value": _short_metric_value(len(deck.slide_ids)), "description": "当前 deck 中已有 slide 数量。"},
                    {"label": "状态", "value": "生成中", "description": "等待补充 SQL 数据后可生成图表页。"},
                ],
            },
        )
    ]


def _add_metric_cards(slide, items: list[dict[str, str]]) -> None:
    cards = items[:5] or [{"label": "关键指标", "value": "-"}]
    gap = 0.12
    x0 = 0.55
    y = 1.18
    total_w = 12.25
    card_w = (total_w - gap * (len(cards) - 1)) / len(cards)
    for index, item in enumerate(cards):
        x = x0 + index * (card_w + gap)
        shape = _add_panel(slide, x, y, card_w, 0.68, color=COLORS["panel_2"])
        frame = shape.text_frame
        frame.clear()
        frame.margin_left = _inches(0.12)
        frame.margin_right = _inches(0.08)
        frame.margin_top = _inches(0.08)
        frame.margin_bottom = _inches(0.04)
        paragraph = frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = item["value"]
        _set_run_font(run, size=16, bold=True, color=COLORS["white"])
        paragraph = frame.add_paragraph()
        run = paragraph.add_run()
        run.text = item["label"]
        _set_run_font(run, size=8.5, color=COLORS["muted"])


def _add_bullets(slide, x: float, y: float, w: float, h: float, title: str, items: list[str]) -> None:
    _add_panel(slide, x, y, w, h)
    _add_section_title(slide, x + 0.15, y + 0.12, title)
    shape = slide.shapes.add_textbox(_inches(x + 0.18), _inches(y + 0.42), _inches(w - 0.34), _inches(h - 0.52))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items[:4] or ["暂无内容"]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.level = 0
        paragraph.text = f"• {_safe_text(item)}"
        for run in paragraph.runs:
            _set_run_font(run, size=9.2, color=COLORS["text"])


def _add_paragraph_panel(slide, x: float, y: float, w: float, h: float, title: str, text: str) -> None:
    _add_panel(slide, x, y, w, h)
    _add_section_title(slide, x + 0.15, y + 0.12, title)
    _add_textbox(slide, x + 0.18, y + 0.45, w - 0.34, h - 0.55, _safe_text(text, "暂无说明"), size=9.2)


def _add_native_chart(slide, rows: list[dict[str, Any]], x: float, y: float, w: float, h: float, chart_type: str) -> bool:
    labels, value_column, values = _chart_fields(rows)
    if not labels or value_column is None or not values:
        return False
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(value_column, values)
    ppt_chart_type = XL_CHART_TYPE.LINE_MARKERS if chart_type == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_shape = slide.shapes.add_chart(
        ppt_chart_type,
        _inches(x),
        _inches(y),
        _inches(w),
        _inches(h),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = f"{value_column} 对比"
    chart.category_axis.tick_labels.font.size = Pt(7)
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    return True


def _add_table(slide, rows: list[dict[str, Any]], x: float, y: float, w: float, h: float, *, max_rows: int) -> bool:
    columns = _columns(rows)
    if not rows or not columns:
        _add_paragraph_panel(slide, x, y, w, h, "数据表", "暂无可导出的表格数据。")
        return False
    visible_rows = [row for row in rows[:max_rows] if isinstance(row, dict)]
    shape = slide.shapes.add_table(
        len(visible_rows) + 1,
        len(columns),
        _inches(x),
        _inches(y),
        _inches(w),
        _inches(h),
    )
    table = shape.table
    for col_index, column in enumerate(columns):
        table.columns[col_index].width = int(_inches(w / len(columns)))
        cell = table.cell(0, col_index)
        _fill_shape(cell, COLORS["panel_2"])
        _set_cell_text(cell, column, size=8.2, bold=True, color=COLORS["white"])
    for row_index, row in enumerate(visible_rows, start=1):
        for col_index, column in enumerate(columns):
            cell = table.cell(row_index, col_index)
            _fill_shape(cell, COLORS["bg"] if row_index % 2 else COLORS["panel"])
            _set_cell_text(cell, _format_cell(row.get(column)), size=7.5, color=COLORS["text"])
    return True


def _add_footer(ppt_slide, draft: SlideDraftDTO, index: int) -> None:
    source = f"Slide {index} | session={draft.session_id} | query={draft.query_id or '-'} | scheme={draft.lineage_summary.get('ppt_scheme_name', '-')}"
    _add_textbox(ppt_slide, 0.55, 7.05, 12.25, 0.22, source, size=7.2, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def _section_labels(draft: SlideDraftDTO) -> dict[str, str]:
    scheme = str(draft.lineage_summary.get("ppt_scheme") or draft.chart_spec.get("ppt_scheme") or "").lower()
    if scheme == "evidence":
        return {"findings": "数据证据", "narrative": "来源链路", "recommendations": "校验动作"}
    if scheme == "dashboard":
        return {"findings": "指标看板", "narrative": "扫读解释", "recommendations": "下钻建议"}
    if scheme == "comparison_page":
        return {"findings": "差异发现", "narrative": "对比解释", "recommendations": "下钻路径"}
    if scheme == "trend_page":
        return {"findings": "趋势发现", "narrative": "波动解释", "recommendations": "验证动作"}
    if scheme == "root_cause_page":
        return {"findings": "可能原因", "narrative": "归因解释", "recommendations": "验证路径"}
    if scheme == "risk_action_page":
        return {"findings": "风险提示", "narrative": "风险解释", "recommendations": "行动建议"}
    if scheme == "summary_page":
        return {"findings": "阶段结论", "narrative": "汇报摘要", "recommendations": "下一步计划"}
    if scheme == "automizer_template":
        return {"findings": "模板插槽", "narrative": "品牌版式映射", "recommendations": "模板检查"}
    if scheme == "banana_iterate":
        return {"findings": "改稿方向", "narrative": "当前版本", "recommendations": "迭代动作"}
    if scheme == "pptagent_plan":
        return {"findings": "规划/编辑/反思", "narrative": "质量自检", "recommendations": "修订动作"}
    if scheme in {"presenton_ai", "presenton_end2end"}:
        return {"findings": "叙事主线", "narrative": "端到端表达", "recommendations": "成稿动作"}
    return {"findings": "核心结论", "narrative": "业务解释", "recommendations": "下一步建议"}


def _add_title(slide, draft: SlideDraftDTO) -> None:
    _add_textbox(slide, 0.55, 0.28, 12.25, 0.38, draft.title, size=20, bold=True, color=COLORS["white"])
    _add_textbox(slide, 0.56, 0.72, 8.95, 0.25, draft.subtitle or "数据分析汇报页", size=9.5, color=COLORS["muted"])
    scheme = _safe_text(draft.lineage_summary.get("ppt_scheme_name") or draft.chart_spec.get("ppt_scheme_name"), "PPT方案")
    _add_textbox(slide, 9.55, 0.72, 3.25, 0.25, f"{scheme} · {draft.page_type}", size=7.8, color=COLORS["cyan"], align=PP_ALIGN.RIGHT)


def _add_slide(prs: Presentation, draft: SlideDraftDTO, index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    _fill_shape(background, COLORS["bg"])
    background.line.fill.background()
    _add_title(slide, draft)
    _add_metric_cards(slide, _metric_items(draft))

    rows = _rows(draft)
    chart_type = str(draft.chart_spec.get("chart_type") or "bar").lower()
    has_chart_data = bool(_chart_fields(rows)[0])
    labels = _section_labels(draft)
    if has_chart_data:
        _add_bullets(slide, 0.55, 2.10, 5.35, 1.35, labels["findings"], draft.findings)
        _add_paragraph_panel(slide, 0.55, 3.60, 5.35, 1.22, labels["narrative"], draft.narrative)
        _add_bullets(slide, 0.55, 5.00, 5.35, 1.32, labels["recommendations"], draft.recommendations)
        _add_section_title(slide, 6.20, 2.03, "Office 原生图表")
        _add_native_chart(slide, rows, 6.15, 2.30, 6.60, 2.15, chart_type)
        _add_section_title(slide, 6.20, 4.62, "可编辑数据表")
        _add_table(slide, rows, 6.15, 4.88, 6.60, 1.88, max_rows=6)
    else:
        _add_bullets(slide, 0.55, 2.10, 5.35, 1.55, labels["findings"], draft.findings)
        _add_paragraph_panel(slide, 0.55, 3.82, 5.35, 1.36, labels["narrative"], draft.narrative)
        _add_bullets(slide, 0.55, 5.35, 5.35, 1.15, labels["recommendations"], draft.recommendations)
        _add_section_title(slide, 6.20, 2.03, "可编辑数据表")
        _add_table(slide, rows, 6.15, 2.32, 6.60, 4.32, max_rows=10)
    _add_footer(slide, draft, index)


def _placeholder_slide(deck: DeckDTO) -> SlideDraftDTO:
    return SlideDraftDTO(
        slide_id="slide_placeholder",
        session_id=deck.session_id,
        page_type="summary",
        title=deck.title,
        subtitle="暂无 slide 内容",
        findings=["当前 deck 还没有 slide，后续可补充后再次导出。"],
        narrative="系统已生成空白 deck 占位导出。",
        recommendations=["补齐 slide 后重新导出。"],
    )


def _write_pptx_with_python_pptx(file_path: Path, deck: DeckDTO, slides: list[SlideDraftDTO]) -> None:
    prs = Presentation()
    prs.slide_width = _inches(13.333)
    prs.slide_height = _inches(7.5)
    for index, slide in enumerate(slides or [_placeholder_slide(deck)], start=1):
        _add_slide(prs, slide, index)
    prs.save(file_path)


def _append_presenton_data_slides(file_path: Path, slides: list[SlideDraftDTO]) -> int:
    if os.getenv("PRESENTON_APPEND_DATA_SLIDES", "0").lower() not in {"1", "true", "yes"}:
        return 0
    data_slides = [slide for slide in slides if _rows(slide)]
    if not data_slides:
        return 0

    prs = Presentation(file_path)
    start_index = len(prs.slides)
    for offset, slide in enumerate(data_slides, start=1):
        appendix_slide = slide.model_copy(
            update={
                "title": f"数据明细：{slide.title}",
                "subtitle": "SQL 结果可编辑数据页",
            }
        )
        _add_slide(prs, appendix_slide, start_index + offset)
    prs.save(file_path)
    return len(data_slides)


def _markdown_table(rows: list[dict[str, Any]], *, max_rows: int = 8, max_columns: int = 6) -> str:
    columns = _columns(rows, max_columns=max_columns)
    if not rows or not columns:
        return ""
    visible_rows = [row for row in rows[:max_rows] if isinstance(row, dict)]
    header = "| " + " | ".join(columns) + " |"
    separator = "| " + " | ".join("---" for _ in columns) + " |"
    body = [
        "| " + " | ".join(_format_cell(row.get(column)).replace("\n", " ")[:80] for column in columns) + " |"
        for row in visible_rows
    ]
    return "\n".join([header, separator, *body])


def _build_presenton_slide_markdown(slide: SlideDraftDTO, index: int) -> str:
    rows = _rows(slide)
    table = _markdown_table(rows)
    sections = [
        f"# {slide.title}",
        slide.subtitle or "",
        "## 核心发现",
        "\n".join(f"- {item}" for item in (slide.findings or ["暂无核心发现"])[:5]),
        "## 业务解释",
        slide.narrative or "暂无业务解释。",
        "## 行动建议",
        "\n".join(f"- {item}" for item in (slide.recommendations or ["继续补充数据后复核。"])[:5]),
        "## 来源",
        f"- session_id: {slide.session_id}",
        f"- query_id: {slide.query_id or '-'}",
        f"- page_type: {slide.page_type}",
    ]
    if table:
        sections.extend(["## SQL 结果预览", table])
    return "\n\n".join(part for part in sections if part)


def _build_presenton_payload(deck: DeckDTO, slides: list[SlideDraftDTO]) -> dict[str, Any]:
    effective_slides = slides or [_placeholder_slide(deck)]
    direct_mode = os.getenv("PRESENTON_DIRECT_DATA_SLIDES", "1").lower() in {"1", "true", "yes"}
    slides_markdown = (
        _build_presenton_direct_slides(deck, effective_slides)
        if direct_mode
        else [
            _build_presenton_slide_markdown(slide, index)
            for index, slide in enumerate(effective_slides, start=1)
        ]
    )
    content = "\n\n---\n\n".join(slides_markdown)
    default_instructions = (
        "请基于给定的 SQL 分析结果生成中文数据分析汇报 PPT。"
        "优先使用 Presenton 模板自身的图表、表格和指标版式。"
        "如果 slide 中包含 PRESENTON_DIRECT_CONTENT_JSON_START，则必须直接采用其中的 chartData、tableData、metrics，不要改写或臆造数据。"
        "保留管理汇报口吻，突出结论、数据证据、来源链路和下一步动作。"
    )
    n_slides = len(slides_markdown) if direct_mode else max(1, len(effective_slides))
    payload = {
        "content": content,
        "instructions": os.getenv("PRESENTON_INSTRUCTIONS", default_instructions),
        "n_slides": max(1, n_slides),
        "language": os.getenv("PRESENTON_LANGUAGE", "Chinese"),
        "template": os.getenv("PRESENTON_TEMPLATE", "general"),
        "include_title_slide": os.getenv("PRESENTON_INCLUDE_TITLE_SLIDE", "false").lower() in {"1", "true", "yes"},
        "include_table_of_contents": False,
        "export_as": "pptx",
        "trigger_webhook": False,
    }
    if direct_mode or os.getenv("PRESENTON_USE_SLIDES_MARKDOWN", "0").lower() in {"1", "true", "yes"}:
        payload["slides_markdown"] = slides_markdown
    return payload


def _presenton_base_url() -> str:
    return os.getenv("PRESENTON_BASE_URL", "http://127.0.0.1:5000").strip().rstrip("/")


def _presenton_timeout() -> int:
    try:
        return int(os.getenv("PRESENTON_TIMEOUT", "300"))
    except ValueError:
        return 300


def _presenton_headers(*, json_body: bool = False) -> dict[str, str]:
    headers: dict[str, str] = {"Accept": "application/json"}
    if json_body:
        headers["Content-Type"] = "application/json"
    authorization = os.getenv("PRESENTON_AUTHORIZATION") or os.getenv("PRESENTON_AUTH_HEADER")
    if authorization:
        headers["Authorization"] = authorization
    cookie = os.getenv("PRESENTON_COOKIE")
    if cookie:
        headers["Cookie"] = cookie
    return headers


def _http_error_detail(exc: HTTPError) -> str:
    try:
        body = exc.read().decode("utf-8", errors="replace")
    except Exception:  # noqa: BLE001
        body = ""
    return f"HTTP {exc.code}: {body or exc.reason}"


def _call_presenton_generate(payload: dict[str, Any]) -> dict[str, Any]:
    base_url = _presenton_base_url()
    if not base_url:
        raise RuntimeError("PRESENTON_BASE_URL is empty")
    endpoint = f"{base_url}/api/v1/ppt/presentation/generate"
    request = Request(
        endpoint,
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers=_presenton_headers(json_body=True),
        method="POST",
    )
    try:
        with urlopen(request, timeout=_presenton_timeout()) as response:
            response_body = response.read().decode("utf-8")
    except HTTPError as exc:
        raise RuntimeError(f"Presenton generate failed: {_http_error_detail(exc)}") from exc
    except URLError as exc:
        raise RuntimeError(f"Presenton service is not reachable at {base_url}: {exc.reason}") from exc
    try:
        result = json.loads(response_body)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"Presenton returned invalid JSON: {response_body[:500]}") from exc
    if not isinstance(result, dict):
        raise RuntimeError("Presenton returned an unexpected response")
    return result


def _presenton_local_candidates(source_path: str) -> list[Path]:
    candidates: list[Path] = []
    raw_path = Path(source_path)
    if raw_path.is_absolute():
        candidates.append(raw_path)
    app_data_dir = os.getenv("PRESENTON_APP_DATA_DIR")
    app_data_roots = [
        Path(app_data_dir) if app_data_dir else None,
        Path("third_party/presenton/app_data"),
        Path("runtime_data/presenton/app_data"),
    ]
    if source_path.startswith("/app_data/"):
        relative = source_path[len("/app_data/") :]
        for root in app_data_roots:
            if root is not None:
                candidates.append(root / relative)
    return candidates


def _copy_or_download_presenton_file(result: dict[str, Any], file_path: Path) -> str:
    source_path = str(result.get("path") or result.get("url") or "").strip()
    if not source_path:
        raise RuntimeError("Presenton response did not include a PPTX path")
    file_path.parent.mkdir(parents=True, exist_ok=True)
    for candidate in _presenton_local_candidates(source_path):
        if candidate.exists() and candidate.is_file():
            shutil.copyfile(candidate, file_path)
            return str(candidate)

    parsed = urlparse(source_path)
    download_url = source_path if parsed.scheme in {"http", "https"} else urljoin(f"{_presenton_base_url()}/", source_path.lstrip("/"))
    request = Request(download_url, headers=_presenton_headers(json_body=False), method="GET")
    try:
        with urlopen(request, timeout=_presenton_timeout()) as response, file_path.open("wb") as output:
            shutil.copyfileobj(response, output)
    except HTTPError as exc:
        raise RuntimeError(f"Presenton PPTX download failed: {_http_error_detail(exc)}") from exc
    except URLError as exc:
        raise RuntimeError(f"Presenton PPTX download failed: {exc.reason}") from exc
    return download_url


def _write_pptx_with_presenton(file_path: Path, deck: DeckDTO, slides: list[SlideDraftDTO]) -> dict[str, str]:
    payload = _build_presenton_payload(deck=deck, slides=slides)
    result = _call_presenton_generate(payload)
    tmp_path = file_path.with_name(f".{file_path.name}.{uuid4().hex}.tmp")
    try:
        source = _copy_or_download_presenton_file(result, tmp_path)
        if not tmp_path.exists() or tmp_path.stat().st_size <= 0:
            raise RuntimeError("Presenton exporter produced an empty file")
        tmp_path.replace(file_path)
        appendix_count = _append_presenton_data_slides(file_path, slides)
    except Exception:
        tmp_path.unlink(missing_ok=True)
        raise
    return {
        "base_url": _presenton_base_url(),
        "presentation_id": str(result.get("presentation_id") or ""),
        "edit_path": str(result.get("edit_path") or ""),
        "source": source,
        "data_appendix_slides": str(appendix_count),
    }


def _write_pptx(file_path: Path, deck: DeckDTO, slides: list[SlideDraftDTO]) -> tuple[str, str | None, dict[str, str]]:
    exporter = os.getenv("VIBE_PPTX_EXPORTER", "presenton").strip().lower()
    if exporter in {"python", "python-pptx", "python_pptx"}:
        _write_pptx_with_python_pptx(file_path, deck=deck, slides=slides)
        return "python-pptx", None, {}

    try:
        metadata = _write_pptx_with_presenton(file_path, deck=deck, slides=slides)
        return "presenton", None, metadata
    except Exception as exc:  # noqa: BLE001
        if os.getenv("VIBE_PRESENTON_STRICT", "0") == "1" or os.getenv("VIBE_PRESENTON_ALLOW_FALLBACK", "0") == "0":
            raise
        _write_pptx_with_python_pptx(file_path, deck=deck, slides=slides)
        return "python-pptx-fallback", str(exc), {"primary": "presenton"}


def export_deck_artifact(deck: DeckDTO, slides: list[SlideDraftDTO] | None = None) -> ArtifactDTO:
    ARTIFACT_ROOT.mkdir(parents=True, exist_ok=True)
    artifact_id = f"artifact_{uuid4().hex[:10]}"
    file_name = f"{deck.deck_id}.pptx"
    file_path = ARTIFACT_ROOT / file_name
    effective_slides = slides or []
    slide_lines: list[str] = []
    for index, slide in enumerate(effective_slides, start=1):
        rows = _rows(slide)
        slide_scheme = str(slide.lineage_summary.get("ppt_scheme", slide.chart_spec.get("ppt_scheme", "presenton_ai")))
        if slide_scheme != "presenton_ai":
            slide_scheme = "presenton_ai"
        slide_lines.extend(
            [
                f"slide_{index}_id={slide.slide_id}",
                f"slide_{index}_title={slide.title}",
                f"slide_{index}_page_type={slide.page_type}",
                f"slide_{index}_scheme={slide_scheme}",
                f"slide_{index}_chart_type={slide.chart_spec.get('chart_type', '')}",
                f"slide_{index}_editable_table={bool(rows)}",
                f"slide_{index}_editable_chart={bool(_chart_fields(rows)[0])}",
                f"slide_{index}_findings={' | '.join(slide.findings)}",
            ]
        )
    exporter_name, exporter_error, exporter_metadata = _write_pptx(file_path, deck=deck, slides=effective_slides)
    package_payload = "\n".join(
        [
            f"deck_id={deck.deck_id}",
            f"session_id={deck.session_id}",
            f"title={deck.title}",
            f"slide_ids={','.join(deck.slide_ids)}",
            f"exporter={exporter_name}",
            f"exporter_primary={exporter_metadata.get('primary', exporter_name)}",
            f"presenton_base_url={exporter_metadata.get('base_url', '')}",
            f"presenton_presentation_id={exporter_metadata.get('presentation_id', '')}",
            f"presenton_edit_path={exporter_metadata.get('edit_path', '')}",
            f"presenton_source={exporter_metadata.get('source', '')}",
            f"presenton_data_appendix_slides={exporter_metadata.get('data_appendix_slides', '0')}",
            f"exporter_fallback={'true' if exporter_error else 'false'}",
            f"exporter_error={exporter_error or ''}",
            "layout=16:9 stable business report",
            *slide_lines,
        ]
    )
    (ARTIFACT_ROOT / f"{deck.deck_id}.md").write_text(package_payload, encoding="utf-8")
    (ARTIFACT_ROOT / f"{deck.deck_id}.html").write_text(
        "<html><head><meta charset='utf-8'><title>"
        + html_escape(deck.title)
        + "</title></head><body><h1>"
        + html_escape(deck.title)
        + "</h1><pre>"
        + html_escape(package_payload)
        + "</pre></body></html>",
        encoding="utf-8",
    )
    return ArtifactDTO(
        artifact_id=artifact_id,
        artifact_type="pptx",
        file_name=file_name,
        download_url=f"/api/v1/artifacts/{artifact_id}/download",
        local_path=str(file_path),
        session_id=deck.session_id,
        deck_id=deck.deck_id,
    )
