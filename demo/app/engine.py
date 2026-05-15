import re
import sqlite3
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Tuple

from .mock_data import ensure_demo_database


FORBIDDEN = ["insert ", "update ", "delete ", "drop ", "alter ", "truncate ", "create "]


@dataclass
class AnalysisResult:
    goal: str
    query_plan: Dict[str, str]
    sql: str
    columns: List[str]
    rows: List[Tuple]
    insight: str
    recommendation: List[str]
    slide: Dict[str, str]


def build_query_plan(goal: str) -> Dict[str, str]:
    text = goal.lower()
    if "上新" in goal or "new" in text:
        return {
            "intent": "新品节奏分析",
            "metric": "new_sku_count",
            "dimension": "platform_name, snapshot_date",
            "window": "最近30天",
            "chart": "line",
        }
    if "价格" in goal or "price" in text:
        return {
            "intent": "价格结构分析",
            "metric": "avg_sale_price, avg_discount",
            "dimension": "brand_name, platform_name",
            "window": "最近30天",
            "chart": "bar",
        }
    return {
        "intent": "商品结构分析",
        "metric": "sku_count",
        "dimension": "category_name, brand_name",
        "window": "最近30天",
        "chart": "bar",
    }


def build_sql(goal: str, plan: Dict[str, str]) -> str:
    text = goal.lower()
    where_category = ""
    if "外套" in goal:
        where_category = " AND c.category_name = '外套'"
    elif "连衣裙" in goal:
        where_category = " AND c.category_name = '连衣裙'"

    if plan["intent"] == "新品节奏分析":
        return f"""
SELECT s.snapshot_date, p.platform_name, COUNT(DISTINCT s.sku_id) AS new_sku_count
FROM product_snapshot s
JOIN dim_platform p ON s.platform_id = p.platform_id
JOIN dim_category c ON s.category_id = c.category_id
WHERE s.snapshot_date >= date('now', '-30 day')
  AND s.is_new = 1
  {where_category}
GROUP BY s.snapshot_date, p.platform_name
ORDER BY s.snapshot_date, new_sku_count DESC
LIMIT 200;
""".strip()

    if "平台" in goal or "platform" in text:
        return f"""
SELECT p.platform_name, b.brand_name,
       ROUND(AVG(s.sale_price), 2) AS avg_sale_price,
       ROUND(AVG(s.sale_price / NULLIF(s.listed_price, 0)), 3) AS avg_discount,
       COUNT(DISTINCT s.sku_id) AS sku_count
FROM product_snapshot s
JOIN dim_brand b ON s.brand_id = b.brand_id
JOIN dim_platform p ON s.platform_id = p.platform_id
JOIN dim_category c ON s.category_id = c.category_id
WHERE s.snapshot_date >= date('now', '-30 day')
  {where_category}
GROUP BY p.platform_name, b.brand_name
ORDER BY avg_sale_price DESC
LIMIT 30;
""".strip()

    return f"""
SELECT b.brand_name, c.category_name, COUNT(DISTINCT s.sku_id) AS sku_count
FROM product_snapshot s
JOIN dim_brand b ON s.brand_id = b.brand_id
JOIN dim_category c ON s.category_id = c.category_id
WHERE s.snapshot_date >= date('now', '-30 day')
  {where_category}
GROUP BY b.brand_name, c.category_name
ORDER BY sku_count DESC
LIMIT 30;
""".strip()


def validate_sql(sql: str) -> None:
    lowered = re.sub(r"\s+", " ", sql.strip().lower())
    if not lowered.startswith("select"):
        raise ValueError("Only SELECT SQL is allowed.")
    if any(token in lowered for token in FORBIDDEN):
        raise ValueError("Dangerous SQL keyword detected.")


def execute_sql(sql: str) -> Tuple[List[str], List[Tuple]]:
    db_path, _ = ensure_demo_database()
    conn = sqlite3.connect(db_path)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()
    cur.execute(sql)
    rows = cur.fetchmany(200)
    columns = [d[0] for d in cur.description] if cur.description else []
    conn.close()
    return columns, [tuple(r) for r in rows]


def summarize(goal: str, plan: Dict[str, str], rows: List[Tuple], columns: List[str]) -> Tuple[str, List[str]]:
    if not rows:
        return "当前条件下没有返回数据，建议放宽筛选条件或扩大时间范围。", [
            "改成最近60天",
            "去掉品类限制",
            "按平台先做总览",
        ]

    first = rows[0]
    lead = ", ".join(f"{columns[i]}={first[i]}" for i in range(min(3, len(columns))))
    insight = (
        f"本轮目标为“{plan['intent']}”。系统在最近30天数据上完成查询，"
        f"共返回 {len(rows)} 行。首条结果：{lead}。"
        "建议下一步做异常点下钻和品牌平台交叉分析。"
    )
    recommendation = [
        "按品牌拆解价格带变化",
        "按平台比较促销深度差异",
        "补一页风险和行动建议",
    ]
    return insight, recommendation


def run_analysis(goal: str) -> AnalysisResult:
    plan = build_query_plan(goal)
    sql = build_sql(goal, plan)
    validate_sql(sql)
    cols, rows = execute_sql(sql)
    insight, recs = summarize(goal, plan, rows, cols)
    slide = {
        "title": f"{plan['intent']} - 自动分析页",
        "summary": insight,
        "chart": plan["chart"],
        "source": "product_snapshot + dimension tables",
        "time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }
    return AnalysisResult(
        goal=goal,
        query_plan=plan,
        sql=sql,
        columns=cols,
        rows=rows,
        insight=insight,
        recommendation=recs,
        slide=slide,
    )


def export_ppt(session_id: str, slides: List[Dict[str, str]]) -> str:
    root = Path(__file__).resolve().parents[2]
    out_dir = root / "output"
    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / f"analysis_{session_id}.pptx"
    try:
        from pptx import Presentation

        prs = Presentation()
        for idx, slide_data in enumerate(slides, start=1):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"{idx}. {slide_data.get('title', 'Analysis Slide')}"
            body = slide.placeholders[1].text_frame
            body.text = slide_data.get("summary", "")
            body.add_paragraph().text = f"Chart Suggestion: {slide_data.get('chart', 'bar')}"
            body.add_paragraph().text = f"Source: {slide_data.get('source', '')}"
            body.add_paragraph().text = f"Generated At: {slide_data.get('time', '')}"
        prs.save(str(path))
        return str(path)
    except Exception:
        fallback = out_dir / f"analysis_{session_id}.md"
        lines = ["# Analysis Report"]
        for idx, s in enumerate(slides, start=1):
            lines.append(f"\n## {idx}. {s.get('title', 'Analysis Slide')}")
            lines.append(s.get("summary", ""))
            lines.append(f"- Chart: {s.get('chart', 'bar')}")
            lines.append(f"- Source: {s.get('source', '')}")
            lines.append(f"- Generated At: {s.get('time', '')}")
        fallback.write_text("\n".join(lines), encoding="utf-8")
        return str(fallback)
